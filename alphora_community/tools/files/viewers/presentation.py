"""
PPT 演示文稿查看器 - 处理 .pptx/.ppt 文件
"""
import os
from typing import Optional, List, Dict, Any, Tuple

from ..utils.common import get_file_info, truncate_text


class PresentationViewer:
    """PPT 演示文稿查看器"""
    
    SUPPORTED_EXTENSIONS = {'.pptx', '.ppt'}
    
    def __init__(self, file_path: str):
        """
        初始化查看器
        
        Args:
            file_path: 文件路径
        """
        self.file_path = file_path
        self.file_info = get_file_info(file_path)
        self.ext = self.file_info['extension']
        
    def view(
        self,
        purpose: str = "preview",
        keyword: Optional[str] = None,
        max_lines: int = 100,
        page_number: Optional[int] = None,
    ) -> str:
        """
        查看 PPT 内容
        
        Args:
            purpose: 查看目的（preview/structure/search）
            keyword: 搜索关键词
            max_lines: 最大返回行数
            page_number: 指定幻灯片页码（从1开始）
            
        Returns:
            格式化的演示文稿内容字符串
        """
        # 智能参数推断
        purpose, warnings = self._infer_and_validate_params(purpose, keyword, page_number)
        
        try:
            from pptx import Presentation
        except ImportError:
            return "❌ 需要安装 python-pptx 库：pip install python-pptx"
        
        try:
            prs = Presentation(self.file_path)
        except Exception as e:
            return f"❌ 无法打开演示文稿: {e}"
        
        total_slides = len(prs.slides)
        
        if purpose == "structure":
            return self._get_structure(prs, total_slides, warnings)
        elif purpose == "search":
            return self._search(prs, keyword, max_lines, warnings)
        elif page_number is not None:
            return self._view_slide(prs, page_number, total_slides, warnings)
        else:  # preview
            return self._preview(prs, total_slides, max_lines, warnings)
    
    def _infer_and_validate_params(
        self,
        purpose: str,
        keyword: Optional[str],
        page_number: Optional[int]
    ) -> Tuple[str, List[str]]:
        """智能推断和校验参数"""
        warnings = []
        
        if keyword and purpose != "search":
            warnings.append(f"⚠️ 检测到 keyword='{keyword}'，已自动切换为 search 模式")
            purpose = "search"
        
        if purpose == "search" and not keyword:
            warnings.append("⚠️ search 模式需要 keyword 参数，已切换为 preview 模式")
            purpose = "preview"
            
        return purpose, warnings
    
    def _get_slide_title(self, slide) -> str:
        """获取幻灯片标题"""
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                return truncate_text(shape.text.strip(), 50)
        return "(无标题)"
    
    def _get_slide_content(self, slide) -> List[str]:
        """获取幻灯片所有文本内容"""
        content = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                content.append(shape.text.strip())
        return content
    
    def _format_header(self, total_slides: int, warnings: List[str]) -> str:
        """格式化文件头信息"""
        lines = [
            f"📊 文件: {self.file_info['name']}",
            f"📦 大小: {self.file_info['size_human']}",
            f"📋 幻灯片数量: {total_slides}",
        ]
        
        if warnings:
            lines.append("")
            for w in warnings:
                lines.append(w)
        
        return '\n'.join(lines)
    
    def _get_structure(self, prs, total_slides: int, warnings: List[str]) -> str:
        """获取演示文稿结构"""
        lines = [
            self._format_header(total_slides, warnings),
            "",
            "【幻灯片列表】"
        ]
        
        for i, slide in enumerate(prs.slides, 1):
            title = self._get_slide_title(slide)
            shape_count = len(slide.shapes)
            
            # 统计内容类型
            text_count = sum(1 for s in slide.shapes if hasattr(s, 'text') and s.text.strip())
            
            lines.append(f"  第{i}页: {title}")
            lines.append(f"        元素: {shape_count}个 (含文本: {text_count}个)")
        
        return '\n'.join(lines)
    
    def _preview(
        self,
        prs,
        total_slides: int,
        max_lines: int,
        warnings: List[str]
    ) -> str:
        """预览演示文稿内容"""
        lines = [
            self._format_header(total_slides, warnings),
            "",
            "【内容预览】"
        ]
        
        content_count = 0
        slides_shown = 0
        
        for i, slide in enumerate(prs.slides, 1):
            if content_count >= max_lines:
                break
            
            title = self._get_slide_title(slide)
            lines.append(f"\n━━━ 第{i}页: {title} ━━━")
            slides_shown = i
            
            for shape in slide.shapes:
                if content_count >= max_lines:
                    break
                if hasattr(shape, 'text') and shape.text.strip():
                    text = shape.text.strip()
                    # 限制单个文本块长度
                    if len(text) > 200:
                        text = text[:200] + "..."
                    lines.append(text)
                    content_count += 1
        
        if slides_shown < total_slides:
            lines.append(f"\n... 还有 {total_slides - slides_shown} 页未显示")
        
        return '\n'.join(lines)
    
    def _view_slide(
        self,
        prs,
        page_number: int,
        total_slides: int,
        warnings: List[str]
    ) -> str:
        """查看指定幻灯片"""
        if page_number < 1 or page_number > total_slides:
            return f"❌ 页码超出范围。该 PPT 共有 {total_slides} 页，请输入 1-{total_slides} 之间的数字"
        
        slide = prs.slides[page_number - 1]
        
        lines = [
            f"📊 文件: {self.file_info['name']}",
            f"📋 第 {page_number}/{total_slides} 页",
        ]
        
        if warnings:
            lines.append("")
            for w in warnings:
                lines.append(w)
        
        lines.append("")
        lines.append("【页面内容】")
        
        content = self._get_slide_content(slide)
        if content:
            for text in content:
                lines.append(text)
                lines.append("")
        else:
            lines.append("(此页没有文本内容)")
        
        # 导航提示
        lines.append("")
        if page_number > 1:
            lines.append(f"💡 上一页: page_number={page_number - 1}")
        if page_number < total_slides:
            lines.append(f"💡 下一页: page_number={page_number + 1}")
        
        return '\n'.join(lines)
    
    def _search(
        self,
        prs,
        keyword: str,
        max_lines: int,
        warnings: List[str]
    ) -> str:
        """搜索演示文稿内容"""
        results = []
        keyword_lower = keyword.lower()
        
        for i, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if hasattr(shape, 'text') and keyword_lower in shape.text.lower():
                    text = shape.text.strip()
                    # 找到关键词上下文
                    idx = text.lower().find(keyword_lower)
                    start = max(0, idx - 30)
                    end = min(len(text), idx + len(keyword) + 50)
                    context = text[start:end]
                    if start > 0:
                        context = "..." + context
                    if end < len(text):
                        context = context + "..."
                    
                    results.append({
                        'page': i,
                        'content': context
                    })
        
        lines = [
            self._format_header(len(prs.slides), warnings),
            "",
            f"🔍 搜索关键词: '{keyword}'",
            f"📋 找到 {len(results)} 处匹配"
        ]
        
        if not results:
            lines.append("")
            lines.append(f"未找到包含 '{keyword}' 的内容")
            return '\n'.join(lines)
        
        lines.append("")
        
        for i, r in enumerate(results[:max_lines], 1):
            lines.append(f"[第{r['page']}页] {r['content']}")
        
        if len(results) > max_lines:
            lines.append(f"\n... 还有 {len(results) - max_lines} 处匹配未显示")
        
        return '\n'.join(lines)
