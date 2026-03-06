"""
PowerPoint Reader

依赖: python-pptx (可选，未安装时给出清晰提示)
"""

import io
from typing import Optional

from alphora.sandbox.tools.inspector.readers import FileContent


def read(data: bytes, path: str, size: int, page: Optional[int] = None, **_kwargs) -> FileContent:
    try:
        from pptx import Presentation
    except ImportError:
        return FileContent(
            text="[Error] python-pptx is required to read PowerPoint files. Install: pip install python-pptx",
            total_lines=0,
            file_type="pptx",
            size=size,
            metadata={"error": "missing_dependency"},
        )

    prs = Presentation(io.BytesIO(data))
    total_slides = len(prs.slides)

    title = ""
    if prs.core_properties and prs.core_properties.title:
        title = prs.core_properties.title

    def extract_slide_text(slide, slide_num: int) -> str:
        lines = [f"── Slide {slide_num} of {total_slides} ──"]
        if slide.shapes.title:
            lines.append(f"Title: {slide.shapes.title.text}")
        lines.append("")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
        return "\n".join(lines)

    if page is not None:
        if page < 1 or page > total_slides:
            return FileContent(
                text=f"[Error] Slide {page} out of range (1-{total_slides})",
                total_lines=0,
                file_type="pptx",
                size=size,
                metadata={"slides": total_slides, "title": title},
            )
        slide = prs.slides[page - 1]
        text = extract_slide_text(slide, page)
    else:
        parts = [extract_slide_text(s, i + 1) for i, s in enumerate(prs.slides)]
        text = "\n\n".join(parts) if parts else "(Empty presentation)"

    lines = text.splitlines()
    return FileContent(
        text=text,
        total_lines=len(lines),
        file_type="pptx",
        size=size,
        metadata={
            "slides": total_slides,
            "title": title,
        },
    )
